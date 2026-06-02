#!/usr/bin/env python3
"""OVLA main system architecture figure -> Final/figs/system_architecture.pptx
Two tight bands: Generation (on-device LLM) L->R, Verification (deterministic, LLM-free) R->L (snake).
Boxes packed with minimal gaps; edit freely in PowerPoint/Keynote afterwards."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

EMU = 914400
def I(v): return Emu(int(v * EMU))

prs = Presentation()
prs.slide_width  = I(13.33)
prs.slide_height = I(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shp = slide.shapes

def C(hex_):
    return RGBColor((hex_>>16)&255, (hex_>>8)&255, hex_&255)

# palette
GEN_BG   = C(0xEAF2FB); VER_BG = C(0xEAF6EE)
LLM_F    = C(0xD6E4F0); LLM_L  = C(0x2E5496)
IR_F     = C(0xFFF2CC); IR_L   = C(0xBF9000)
GATE_F   = C(0xFCE4D6); GATE_L = C(0xC55A11)
USER_F   = C(0xE2D4F0); USER_L = C(0x7030A0)
VER_F    = C(0xD9EAD3); VER_L  = C(0x538135)
DEP_F    = C(0xC6E0B4); DEP_L  = C(0x375623)
REJ_F    = C(0xF8CBAD); REJ_L  = C(0xC00000)
NL_F     = C(0xF2F2F2); NL_L   = C(0x7F7F7F)
TXT      = C(0x1F1F1F)
GRAY     = C(0x595959); ORANGE = C(0xC55A11); RED = C(0xC00000)

boxes = {}  # name -> (x,y,w,h)

def box(name, x, y, w, h, title, sub=None, fill=LLM_F, line=LLM_L,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, tcolor=TXT, tsize=11, ssize=8, bold=True):
    s = shp.add_shape(shape, I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(1.25)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = title.split("\n")
    for li, ln in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.size = Pt(tsize); r.font.bold = bold; r.font.color.rgb = tcolor
        r.font.name = "Arial"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(ssize); r2.font.bold = False; r2.font.color.rgb = tcolor
        r2.font.italic = True; r2.font.name = "Arial"
    boxes[name] = (x, y, w, h)
    return s

def band(x, y, w, h, label, fill):
    s = shp.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill; s.shadow.inherit = False
    tf = s.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Pt(2); tf.margin_left = Pt(6)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = C(0x404040)
    r.font.name = "Arial"
    # send to back
    sp = s._element; sp.getparent().remove(sp); shp._spTree.insert(2, sp)
    return s

def ctr(name): x,y,w,h = boxes[name]; return (x+w/2, y+h/2)
def edge(name, side):
    x,y,w,h = boxes[name]
    return {'l':(x,y+h/2),'r':(x+w,y+h/2),'t':(x+w/2,y),'b':(x+w/2,y+h)}[side]

def style_ln(conn, color, width=1.5, dash=None):
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash}); ln.append(d)
    te = ln.makeelement(qn('a:tailEnd'), {'type':'triangle','w':'med','len':'med'})
    ln.append(te)

def arrow(p1, p2, color=GRAY, width=1.5, dash=None):
    c = shp.add_connector(MSO_CONNECTOR.STRAIGHT, I(p1[0]), I(p1[1]), I(p2[0]), I(p2[1]))
    style_ln(c, color, width, dash)
    return c

def label(x, y, text, color=GRAY, size=8, w=1.4):
    s = shp.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(0.22))
    s.fill.background(); s.line.fill.background(); s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = False
    tf.margin_left=Pt(0); tf.margin_right=Pt(0); tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name="Arial"

# ---------------- background bands ----------------
band(0.20, 0.78, 12.95, 1.95, "Phase 1   Generation  —  on-device LLM (≤9B, 4-bit, no cloud)", GEN_BG)
band(0.20, 4.60, 12.95, 2.55, "Phase 2   Verification  —  deterministic, LLM-free  (pass → deploy / fail → repair or reject, fail-closed)", VER_BG)

# ---------------- Phase 1 band (L -> R) ----------------
y1 = 1.25; h1 = 1.05
xw = 1.50; g = 0.114; x0 = 0.30
def px(i): return x0 + i*(xw+g)

box("nl",     px(0), y1, xw, h1, "NL command", "end-user", fill=NL_F, line=NL_L, shape=MSO_SHAPE.ROUNDED_RECTANGLE, tsize=10.5)
box("intent", px(1), y1, xw, h1, "Intent\nAnalysis", "decompose / plan", tsize=10.5)
box("svc",    px(2), y1, xw, h1, "Service\nMapping", "+ arg / enum resolve", tsize=10.5)
box("dev",    px(3), y1, xw, h1, "Device\nMapping", "selector", tsize=10.5)
box("ir",     px(4), y1, xw, h1, "Timeline IR", "extract", fill=IR_F, line=IR_L, tsize=10.5)
box("feas",   px(5), y1, xw, h1, "Feasibility\nGate", "IR ∈ L(G)", fill=GATE_F, line=GATE_L, shape=MSO_SHAPE.HEXAGON, tsize=10)
box("conf",   px(6), y1, xw, h1, "User\nConfirm", "plain-language\nrendering", fill=USER_F, line=USER_L, tsize=10.5)
box("joi",    px(7), y1, xw, h1, "JoI\nGeneration", "lowering", tsize=10.5)

# ---------------- Phase 2 band (R -> L, snake) ----------------
y2 = 5.35; h2 = 1.05
b_w = 2.00; b_g = 0.15
right_edge = px(7) + xw            # align L1 under JoI gen (right)
def bx(i_from_right, w=b_w):       # i_from_right=0 is rightmost
    return right_edge - w - i_from_right*(b_w+b_g)

box("l1",     bx(0), y2, b_w, h2, "L1 Static Check", "JoI well-formedness", fill=VER_F, line=VER_L, tsize=11)
box("fsm",    bx(1), y2, b_w, h2, "IR → FSM\n+ Event Synthesis", "boundary events", fill=VER_F, line=VER_L, tsize=11)
box("sims",   bx(2), y2, b_w, h2, "IR sim  ‖  JoI sim", "execute", fill=VER_F, line=VER_L, tsize=11)
box("teq",    bx(3), y2, 1.85, h2, "Trace-\nEquivalence", "±tolerance", fill=GATE_F, line=GATE_L, shape=MSO_SHAPE.HEXAGON, tsize=10.5)
box("dep",    bx(3)-1.65-b_g, y2, 1.55, h2, "Deploy ✓", "to hub", fill=DEP_F, line=DEP_L, tsize=12)

# Reject in the gap, far left
box("rej",    0.30, 3.10, 1.25, 0.95, "Reject", "fail-closed", fill=REJ_F, line=REJ_L, tsize=11)

# ---------------- forward connectors: Phase 1 ----------------
chain1 = ["nl","intent","svc","dev","ir","feas","conf","joi"]
for a,b in zip(chain1, chain1[1:]):
    arrow(edge(a,'r'), edge(b,'l'), GRAY, 1.5)

# JoI gen -> L1 (handoff, JoI code) : short vertical
arrow(edge("joi",'b'), edge("l1",'t'), GRAY, 1.75)
label(edge("joi",'b')[0]-0.5, (edge("joi",'b')[1]+edge("l1",'t')[1])/2-0.12, "JoI code", GRAY, 8, w=1.0)

# confirmed IR -> FSM box
arrow(edge("conf",'b'), edge("fsm",'t'), GRAY, 1.5)
label(edge("conf",'b')[0]-0.6, edge("conf",'b')[1]+0.5, "confirmed IR", GRAY, 8, w=1.3)

# ---------------- forward connectors: Phase 2 (R -> L) ----------------
chain2 = ["l1","fsm","sims","teq","dep"]
for a,b in zip(chain2, chain2[1:]):
    arrow(edge(a,'l'), edge(b,'r'), GRAY, 1.6)
label((ctr("teq")[0]+ctr("dep")[0])/2-0.5, y2-0.30, "pass", DEP_L, 8, w=1.0)

# ---------------- repair loop: trace-eq -> JoI gen ----------------
arrow(edge("teq",'t'), edge("joi",'b'), ORANGE, 1.75, dash='dash')
label((edge("teq",'t')[0]+edge("joi",'b')[0])/2-0.7, 3.55, "repair: regenerate JoI", ORANGE, 8.5, w=2.6)
label((edge("teq",'t')[0]+edge("joi",'b')[0])/2-0.7, 3.78, "(counterexample)", ORANGE, 8, w=2.6)

# ---------------- reject arrows ----------------
arrow(edge("feas",'b'), edge("rej",'r'), RED, 1.5, dash='dash')        # feasibility fail
arrow(edge("teq",'l'),  edge("rej",'b'), RED, 1.5, dash='dash')        # trace-eq irreparable
label(edge("feas",'b')[0]-1.1, edge("feas",'b')[1]+0.05, "fail", RED, 8, w=0.8)

prs.save("Final/figs/system_architecture.pptx")
print("saved Final/figs/system_architecture.pptx")
